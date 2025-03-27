from PyQt6.QtWidgets import (QMainWindow, QVBoxLayout, QHBoxLayout, QLabel,
                             QPushButton, QTextEdit, QScrollArea, QWidget)
from PyQt6.QtCore import Qt
from PyQt6.QtGui import QPixmap, QImage
import os
from PIL import Image
import fitz  # PyMuPDF
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import io

class PreviewWindow(QMainWindow):
    def __init__(self, file_path: str, parent=None):
        super().__init__(parent)
        self.file_path = file_path
        self.setWindowTitle(f"Preview - {file_path}")
        self.setMinimumSize(800, 600)
        
        # Create central widget and layout
        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        layout = QVBoxLayout(central_widget)
        
        # Header with file name
        header = QHBoxLayout()
        file_label = QLabel(f"File: {self.file_path}")
        file_label.setStyleSheet("font-size: 14px; font-weight: bold; color: black;")
        header.addWidget(file_label)
        layout.addLayout(header)
        
        # Preview area
        self.preview_area = QScrollArea()
        self.preview_area.setWidgetResizable(True)
        self.preview_area.setStyleSheet("""
            QScrollArea {
                border: 1px solid #ccc;
                border-radius: 5px;
                background-color: white;
            }
        """)
        layout.addWidget(self.preview_area)
        
        # Content widget
        self.content_widget = QWidget()
        self.content_layout = QVBoxLayout(self.content_widget)
        self.preview_area.setWidget(self.content_widget)
        
        # Close button
        close_button = QPushButton("Close")
        close_button.setStyleSheet("""
            QPushButton {
                background-color: #f44336;
                color: white;
                padding: 8px 16px;
                border: none;
                border-radius: 4px;
            }
            QPushButton:hover {
                background-color: #d32f2f;
            }
        """)
        close_button.clicked.connect(self.close)
        layout.addWidget(close_button)
        
        # Load preview after UI setup
        self.load_preview()

    def load_preview(self):
        try:
            file_ext = self.file_path.lower().split('.')[-1]
            
            if file_ext in ['jpg', 'jpeg', 'png', 'gif', 'bmp']:
                self.preview_image()
            elif file_ext == 'pdf':
                self.preview_pdf()
            elif file_ext == 'docx':
                self.preview_docx()
            elif file_ext == 'xlsx':
                self.preview_xlsx()
            elif file_ext == 'pptx':
                self.preview_pptx()
            elif file_ext == 'txt':
                self.preview_text()
            else:
                self.show_error("Unsupported file type")
                
        except Exception as e:
            self.show_error(f"Error loading preview: {str(e)}")

    def preview_image(self):
        try:
            image = Image.open(self.file_path)
            # Resize image to fit window while maintaining aspect ratio
            preview_size = self.preview_area.size()
            image.thumbnail(preview_size, Image.Resampling.LANCZOS)
            
            # Convert to QPixmap
            buffer = io.BytesIO()
            image.save(buffer, format='PNG')
            qimage = QImage.fromData(buffer.getvalue())
            pixmap = QPixmap.fromImage(qimage)
            
            # Create label and set pixmap
            label = QLabel()
            label.setPixmap(pixmap)
            label.setAlignment(Qt.AlignmentFlag.AlignCenter)
            self.content_layout.addWidget(label)
            
        except Exception as e:
            self.show_error(f"Error loading image: {str(e)}")

    def preview_pdf(self):
        try:
            doc = fitz.open(self.file_path)
            page = doc[0]
            pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
            
            # Convert to QPixmap
            img = QImage(pix.samples, pix.width, pix.height, pix.stride, QImage.Format.Format_RGB888)
            pixmap = QPixmap.fromImage(img)
            
            # Create label and set pixmap
            label = QLabel()
            label.setPixmap(pixmap)
            label.setAlignment(Qt.AlignmentFlag.AlignCenter)
            self.content_layout.addWidget(label)
            
            doc.close()
            
        except Exception as e:
            self.show_error(f"Error loading PDF: {str(e)}")

    def preview_docx(self):
        try:
            doc = Document(self.file_path)
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            
            text_edit = QTextEdit()
            text_edit.setPlainText(text)
            text_edit.setReadOnly(True)
            text_edit.setStyleSheet("""
                QTextEdit {
                    color: black;
                    background-color: white;
                    border: none;
                }
            """)
            self.content_layout.addWidget(text_edit)
            
        except Exception as e:
            self.show_error(f"Error loading DOCX: {str(e)}")

    def preview_xlsx(self):
        try:
            wb = load_workbook(self.file_path, data_only=True)
            sheet = wb.active
            
            text = ""
            for row in sheet.rows:
                row_text = "\t".join(str(cell.value) if cell.value is not None else "" for cell in row)
                text += row_text + "\n"
            
            text_edit = QTextEdit()
            text_edit.setPlainText(text)
            text_edit.setReadOnly(True)
            text_edit.setStyleSheet("""
                QTextEdit {
                    color: black;
                    background-color: white;
                    border: none;
                }
            """)
            self.content_layout.addWidget(text_edit)
            
        except Exception as e:
            self.show_error(f"Error loading XLSX: {str(e)}")

    def preview_pptx(self):
        try:
            prs = Presentation(self.file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n---\n\n"
            
            text_edit = QTextEdit()
            text_edit.setPlainText(text)
            text_edit.setReadOnly(True)
            text_edit.setStyleSheet("""
                QTextEdit {
                    color: black;
                    background-color: white;
                    border: none;
                }
            """)
            self.content_layout.addWidget(text_edit)
            
        except Exception as e:
            self.show_error(f"Error loading PPTX: {str(e)}")

    def preview_text(self):
        try:
            with open(self.file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            text_edit = QTextEdit()
            text_edit.setPlainText(text)
            text_edit.setReadOnly(True)
            text_edit.setStyleSheet("""
                QTextEdit {
                    color: black;
                    background-color: white;
                    border: none;
                }
            """)
            self.content_layout.addWidget(text_edit)
            
        except Exception as e:
            self.show_error(f"Error loading text file: {str(e)}")

    def show_error(self, message: str):
        error_label = QLabel(message)
        error_label.setStyleSheet("color: red; font-weight: bold;")
        error_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.content_layout.addWidget(error_label) 